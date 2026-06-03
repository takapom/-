from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("uav_harness_progress_architecture.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Hiragino Sans"
FONT_FALLBACK = "Yu Gothic"

BG = RGBColor(248, 250, 252)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
FAINT = RGBColor(226, 232, 240)
TEAL = RGBColor(15, 118, 110)
BLUE = RGBColor(37, 99, 235)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(190, 18, 60)
WHITE = RGBColor(255, 255, 255)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_title,
        slide_research_position,
        slide_progress_summary,
        slide_architecture,
        slide_nl_to_ir,
        slide_place_resolution,
        slide_sitl_gate,
        slide_evaluation_design,
        slide_fact_status,
        slide_next_steps,
    ]
    for idx, builder in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        paint_background(slide)
        builder(slide)
        add_footer(slide, idx, len(slides))

    prs.save(OUT)


def paint_background(slide) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()


def add_footer(slide, idx: int, total: int) -> None:
    add_text(
        slide,
        "Natural-Language UAV Mission Harness",
        Inches(0.55),
        Inches(7.02),
        Inches(5.3),
        Inches(0.25),
        size=8.5,
        color=MUTED,
    )
    add_text(
        slide,
        f"{idx:02d} / {total:02d}",
        Inches(11.9),
        Inches(7.02),
        Inches(0.8),
        Inches(0.25),
        size=8.5,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, Inches(0.65), Inches(0.42), Inches(11.7), Inches(0.55), size=25, bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.68), Inches(1.02), Inches(10.8), Inches(0.35), size=10.5, color=MUTED)


def add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    size: float = 14,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    try:
        run.font._element.rPr.rFonts.set("eastAsia", FONT)
    except AttributeError:
        run.font.name = FONT_FALLBACK


def add_bullets(slide, items: list[str], x, y, w, h, *, size: float = 13, color: RGBColor = TEXT) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color


def card(slide, x, y, w, h, title: str, body: str, *, accent: RGBColor = TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = FAINT
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, title, x + Inches(0.25), y + Inches(0.18), w - Inches(0.45), Inches(0.28), size=11.5, bold=True)
    add_text(slide, body, x + Inches(0.25), y + Inches(0.55), w - Inches(0.45), h - Inches(0.72), size=10.3, color=MUTED)


def node(slide, text: str, x, y, w, h, *, fill: RGBColor = WHITE, line: RGBColor = FAINT, size: float = 10.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = TEXT
    return shape


def arrow(slide, x1, y1, x2, y2, *, color: RGBColor = MUTED) -> None:
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.4)
    line.line.end_arrowhead = True


def pill(slide, text: str, x, y, w, *, fill: RGBColor = TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE


def slide_title(slide) -> None:
    add_text(slide, "自然言語UAVミッション変換ハーネス", Inches(0.78), Inches(1.35), Inches(10.8), Inches(0.72), size=30, bold=True)
    add_text(slide, "進捗サマリと現行アーキテクチャ", Inches(0.82), Inches(2.12), Inches(8.2), Inches(0.45), size=17, color=TEAL, bold=True)
    add_text(slide, "Mission IR / Validator / Emitter / SITL Gate / 評価ログ", Inches(0.84), Inches(2.75), Inches(8.9), Inches(0.35), size=11.5, color=MUTED)
    for i, label in enumerate(["Natural Language", "Mission IR", "SITL Gate", "Evaluation"]):
        pill(slide, label, Inches(0.84 + i * 2.1), Inches(3.55), Inches(1.72), fill=[TEAL, BLUE, AMBER, RED][i])
    card(slide, Inches(8.25), Inches(4.7), Inches(3.8), Inches(1.15), "Current Fact", "目的地command生成は成功。SITL到達成功は未確認。現状態ではPreflightが実行を止める。", accent=AMBER)


def slide_research_position(slide) -> None:
    add_title(slide, "研究の位置づけ", "LLMで直接ドローンを飛ばすのではなく、安全制約付き変換ハーネスを研究対象にする。")
    card(slide, Inches(0.75), Inches(1.65), Inches(3.75), Inches(1.55), "対象", "自然言語で与えられるUAVミッション生成・変更を、実行可能なMission IRへ変換する。", accent=TEAL)
    card(slide, Inches(4.75), Inches(1.65), Inches(3.75), Inches(1.55), "境界", "LLMは高レベル仕様・修復案に限定。低レベル制御はArduPilot/Pixhawkに残す。", accent=BLUE)
    card(slide, Inches(8.75), Inches(1.65), Inches(3.75), Inches(1.55), "評価", "SITL、Validator、Runtime Monitor、評価ログで安全性と再現性を測る。", accent=AMBER)
    add_bullets(
        slide,
        [
            "主張: Mission IR + 決定的Validator + SITL Gateにより、直接コード生成より安全な変換経路を作る。",
            "初期比較: C0 Direct Code / C2 Mission IR / C6 Patch Harness。",
            "現段階: 実機ではなくArduPilot SITL上の再現可能な評価に集中。",
        ],
        Inches(1.05),
        Inches(4.05),
        Inches(10.8),
        Inches(1.6),
        size=13,
    )


def slide_progress_summary(slide) -> None:
    add_title(slide, "ここまでの実装進捗", "自然言語からMission IR、MAVLink command、評価ログまでの骨格が動作する。")
    items = [
        ("Mission IR", "Pydantic schema / local NED action / constraints / patch"),
        ("Validator", "高度・速度・geofence・no-fly・危険指示を決定的に判定"),
        ("Emitter", "検証済みIRをMAVLink commandへ固定変換"),
        ("NL Generator", "日本語地点指示、inline local offset、外部geo opt-in"),
        ("SITL Gate", "Preflight / segmented goto / progress monitor / reset-required"),
        ("Evaluation", "C0/C2/C6比較、JSONLログ、pytest 61 passed"),
    ]
    for idx, (t, b) in enumerate(items):
        col = idx % 2
        row = idx // 2
        card(slide, Inches(0.75 + col * 6.05), Inches(1.45 + row * 1.65), Inches(5.55), Inches(1.18), t, b, accent=[TEAL, BLUE, AMBER, RED, TEAL, BLUE][idx])


def slide_architecture(slide) -> None:
    add_title(slide, "現行アーキテクチャ", "人間の自然言語を、検証可能なMission IRと決定的なSITL commandへ落とす。")
    labels = [
        "日本語タスク",
        "NL Parser\n/ Geo Resolver",
        "Mission IR\n/ Patch",
        "Validator",
        "Emitter",
        "SITL Gate",
        "MAVLink\nArduPilot",
        "Telemetry\n/ Eval Log",
    ]
    xs = [0.7, 2.15, 3.8, 5.35, 6.75, 8.15, 9.75, 11.0]
    widths = [1.15, 1.3, 1.15, 1.05, 1.05, 1.15, 1.05, 1.25]
    for i, label in enumerate(labels):
        node(slide, label, Inches(xs[i]), Inches(2.05), Inches(widths[i]), Inches(0.86), fill=WHITE, line=[TEAL, TEAL, BLUE, BLUE, AMBER, AMBER, RED, RED][i], size=8.7)
        if i < len(labels) - 1:
            arrow(slide, Inches(xs[i] + widths[i] + 0.05), Inches(2.48), Inches(xs[i + 1] - 0.05), Inches(2.48))
    card(slide, Inches(0.85), Inches(4.0), Inches(3.55), Inches(1.35), "LLM / Codexの責務", "自然言語から高レベル仕様を作る。安全判定や低レベル制御の最終責任は持たない。", accent=TEAL)
    card(slide, Inches(4.8), Inches(4.0), Inches(3.55), Inches(1.35), "Harnessの責務", "型・制約・状態・実行前後の検査で、実行してよいcommandだけを通す。", accent=BLUE)
    card(slide, Inches(8.75), Inches(4.0), Inches(3.55), Inches(1.35), "ArduPilotの責務", "姿勢制御、安定化、MAVLink commandの実行、SITLテレメトリを担う。", accent=AMBER)


def slide_nl_to_ir(slide) -> None:
    add_title(slide, "自然言語からMission IRへ", "Mission IRは自然言語そのものではなく、実行前に検証できる中間表現。")
    node(slide, "「A地点まで行って」", Inches(0.9), Inches(1.75), Inches(2.15), Inches(0.72), line=TEAL)
    arrow(slide, Inches(3.05), Inches(2.11), Inches(3.75), Inches(2.11))
    node(slide, "意図・地点・高度\n待機・着陸", Inches(3.85), Inches(1.75), Inches(1.85), Inches(0.72), line=TEAL)
    arrow(slide, Inches(5.7), Inches(2.11), Inches(6.4), Inches(2.11))
    node(slide, "Mission IR\nJSON/Pydantic", Inches(6.5), Inches(1.75), Inches(2.0), Inches(0.72), line=BLUE)
    arrow(slide, Inches(8.5), Inches(2.11), Inches(9.2), Inches(2.11))
    node(slide, "Validator\naccept/reject", Inches(9.3), Inches(1.75), Inches(2.0), Inches(0.72), line=AMBER)
    add_bullets(
        slide,
        [
            "IRに含めるもの: actions, local NED position, constraints, emergency_policy。",
            "IRに含めないもの: 毎回生成される任意Pythonコード、低レベル姿勢制御。",
            "Validatorにより、曖昧・危険・制約違反の候補はSITLへ送らない。",
        ],
        Inches(1.05),
        Inches(3.65),
        Inches(10.6),
        Inches(1.55),
        size=13,
    )


def slide_place_resolution(slide) -> None:
    add_title(slide, "地名目的地の解決フロー", "辞書内地点・inline local offset・外部geoの3経路を持つ。")
    node(slide, "自然言語\n目的地", Inches(0.8), Inches(2.0), Inches(1.35), Inches(0.8), line=TEAL)
    arrow(slide, Inches(2.15), Inches(2.4), Inches(2.7), Inches(2.4))
    node(slide, "local_places\n辞書", Inches(2.8), Inches(1.25), Inches(1.45), Inches(0.75), line=BLUE)
    node(slide, "inline NED\n北/東/高度", Inches(2.8), Inches(2.25), Inches(1.45), Inches(0.75), line=BLUE)
    node(slide, "Nominatim\n/ Overpass", Inches(2.8), Inches(3.25), Inches(1.45), Inches(0.75), line=BLUE)
    arrow(slide, Inches(4.25), Inches(2.4), Inches(5.2), Inches(2.4))
    node(slide, "lat/lon\nor local NED", Inches(5.3), Inches(2.0), Inches(1.45), Inches(0.8), line=AMBER)
    arrow(slide, Inches(6.75), Inches(2.4), Inches(7.5), Inches(2.4))
    node(slide, "SITL home\nrelative NED", Inches(7.6), Inches(2.0), Inches(1.55), Inches(0.8), line=AMBER)
    arrow(slide, Inches(9.15), Inches(2.4), Inches(9.9), Inches(2.4))
    node(slide, "Mission IR\nposition", Inches(10.0), Inches(2.0), Inches(1.55), Inches(0.8), line=TEAL)
    card(slide, Inches(1.0), Inches(5.1), Inches(5.45), Inches(1.0), "実装済みの拒否条件", "unknown / ambiguous / home未指定 / sensitive facility / geo error", accent=RED)
    card(slide, Inches(6.85), Inches(5.1), Inches(4.9), Inches(1.0), "今回の確認例", "Alexander Maconochie Centre -> 約1.08km先のlocal NED目標", accent=AMBER)


def slide_sitl_gate(slide) -> None:
    add_title(slide, "SITL実行ガード", "目的地到達以前に、SITL状態が実行可能かを明示的に確認する。")
    rows = [
        ("1", "Preflight Gate", "armed / landed / altitude / speed / batteryを確認"),
        ("2", "Segmented Route", "長距離goto_local_nedを約100m単位へ分割"),
        ("3", "Progress Monitor", "距離が縮まらない場合はstuckとして停止"),
        ("4", "Reset Required", "失敗後はnext_required_action: sitl_resetを記録"),
    ]
    for i, (n, t, b) in enumerate(rows):
        y = Inches(1.42 + i * 1.08)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y, Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = [TEAL, BLUE, AMBER, RED][i]
        circle.line.fill.background()
        add_text(slide, n, Inches(0.95), y + Inches(0.09), Inches(0.42), Inches(0.2), size=8.8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, t, Inches(1.6), y - Inches(0.02), Inches(3.1), Inches(0.3), size=14.5, bold=True)
        add_text(slide, b, Inches(4.4), y, Inches(7.3), Inches(0.34), size=12, color=MUTED)
    card(slide, Inches(1.1), Inches(6.05), Inches(10.8), Inches(0.72), "現ファクト", "14550ではheartbeatあり。ただしSITL状態がIN_AIRのためPreflightで止め、command_countは0。", accent=AMBER)


def slide_evaluation_design(slide) -> None:
    add_title(slide, "評価設計: C0 / C2 / C6", "提案方式を一括評価せず、構成要素の効果を分離する。")
    card(slide, Inches(0.85), Inches(1.55), Inches(3.65), Inches(2.1), "C0 Direct Code", "LLM由来コード生成の比較条件。現ハーネスでは直接実行を無効化し、ログのみ残す。", accent=RED)
    card(slide, Inches(4.75), Inches(1.55), Inches(3.65), Inches(2.1), "C2 Mission IR", "自然言語タスクを型付きIRへ変換し、ValidatorとEmitterを通して実行候補にする。", accent=TEAL)
    card(slide, Inches(8.65), Inches(1.55), Inches(3.65), Inches(2.1), "C6 Patch Harness", "safe boundaryまで現ミッションを維持し、検証済みpatchを状態に適用する。", accent=BLUE)
    add_bullets(
        slide,
        [
            "評価ログ: task_id, condition_id, Mission IR, patch, validator_result, emitted_commands, sitl_result, telemetry_summary。",
            "主指標候補: 安全制約違反率、危険指示誤受理率、変更処理成功率、状態整合率。",
            "現状テスト: pytest 61 passed / ruff passed / dry-runでcommand分割を確認。",
        ],
        Inches(1.0),
        Inches(4.45),
        Inches(10.9),
        Inches(1.55),
        size=12.7,
    )


def slide_fact_status(slide) -> None:
    add_title(slide, "確認済みファクトと未確認事項", "到達成功はまだ主張しない。生成成功とSITL実行成功を分離する。")
    card(slide, Inches(0.8), Inches(1.45), Inches(5.7), Inches(2.0), "確認済み", "自然言語 -> 地名解決 -> Mission IR -> Validator -> 11分割command生成。約1.08km先の目標を約97.8mごとに分割。", accent=TEAL)
    card(slide, Inches(6.85), Inches(1.45), Inches(5.2), Inches(2.0), "未確認", "修正後の実装でSITLが目的地へ到達した実績はまだない。現在のSITLはPreflightで停止。", accent=AMBER)
    node(slide, "14552\nheartbeatなし", Inches(1.3), Inches(4.55), Inches(2.0), Inches(0.85), line=RED)
    node(slide, "14551\nheartbeatなし", Inches(4.0), Inches(4.55), Inches(2.0), Inches(0.85), line=RED)
    node(slide, "14550\nheartbeatあり / IN_AIR", Inches(6.7), Inches(4.55), Inches(2.25), Inches(0.85), line=AMBER)
    node(slide, "command_count: 0\nrequires_sitl_reset", Inches(9.65), Inches(4.55), Inches(2.25), Inches(0.85), line=AMBER)


def slide_next_steps(slide) -> None:
    add_title(slide, "次の作業", "SITLをリセットし、実到達ログを取って研究評価へ接続する。")
    card(slide, Inches(0.9), Inches(1.35), Inches(3.65), Inches(1.35), "1. SITL Reset", "地上・非armed・正常battery/failsafe状態へ戻す。QGroundControlでも確認。", accent=TEAL)
    card(slide, Inches(4.85), Inches(1.35), Inches(3.65), Inches(1.35), "2. Short Mission", "A地点など短距離でarm/takeoff/goto/landの健全性を確認。", accent=BLUE)
    card(slide, Inches(8.8), Inches(1.35), Inches(3.65), Inches(1.35), "3. External Geo", "Alexander Maconochie Centreの長距離分割実行をログ化。", accent=AMBER)
    add_bullets(
        slide,
        [
            "成功時: segmentごとの到達、最終NED誤差、land完了、telemetry_summaryを評価ログに保存。",
            "失敗時: failed_command、preflight、progress stall、reset_requiredを失敗分類として残す。",
            "論文化観点: semantic success と execution success を分離して、ハーネス設計条件として示す。",
        ],
        Inches(1.05),
        Inches(4.1),
        Inches(10.6),
        Inches(1.55),
        size=13,
    )


if __name__ == "__main__":
    main()
